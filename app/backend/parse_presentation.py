import json
import os
from typing import List, Dict, Any
import tiktoken
from pptx import Presentation

def count_tokens(text: str) -> int:
    """Count the number of tokens in a text string."""
    encoding = tiktoken.get_encoding("cl100k_base")
    return len(encoding.encode(text))

def extract_text_from_slide(slide):
    """Extract text from a slide."""
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"
    return text

def chunk_presentation(prs: Presentation, max_tokens: int = 500) -> List[Dict[str, Any]]:
    """Split presentation into chunks that fit within token limit."""
    chunks = []
    current_chunk = []
    current_tokens = 0
    
    for i, slide in enumerate(prs.slides):
        slide_text = extract_text_from_slide(slide)
        slide_dict = {
            "slide_number": i + 1,
            "content": slide_text
        }
        slide_json = json.dumps(slide_dict)
        slide_tokens = count_tokens(slide_json)
        
        if current_tokens + slide_tokens > max_tokens and current_chunk:
            chunks.append(current_chunk)
            current_chunk = []
            current_tokens = 0
        
        current_chunk.append(slide_dict)
        current_tokens += slide_tokens
    
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks

def parse_pptx(file_path: str, max_tokens: int = 500) -> List[Dict[str, Any]]:
    """Parse PPTX file and return chunks of data within token limit."""
    prs = Presentation(file_path)
    return chunk_presentation(prs, max_tokens)

def process_presentation_files(data_folder: str = "../../data"):
    """Process all PPT and PPTX files in the data folder."""
    for file_name in os.listdir(data_folder):
        file_path = os.path.join(data_folder, file_name)
        
        if not os.path.isfile(file_path):
            continue
            
        file_ext = os.path.splitext(file_name)[1].lower()
        
        try:
            if file_ext in ['.pptx', '.ppt']:
                chunks = parse_pptx(file_path)
                
                # Save chunks as JSON files
                for i, chunk in enumerate(chunks):
                    output_path = f"{file_path}_chunk_{i}.json"
                    with open(output_path, 'w', encoding='utf-8') as f:
                        json.dump(chunk, f, indent=2)
                        
                print(f"Successfully processed: {file_name}")
                
        except Exception as e:
            print(f"Error processing {file_name}: {str(e)}")

if __name__ == "__main__":
    process_presentation_files() 